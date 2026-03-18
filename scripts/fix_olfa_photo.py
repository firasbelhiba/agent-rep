import sys, os
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches

prs = Presentation('AgentRep-PitchDeck.pptx')
slide = prs.slides[1]

# Add olfa photo back with proper sizing (portrait aspect ratio)
photo_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'olfa.png')
photo_left = Inches(5.2) + Inches(0.3)
photo_top = Inches(1.0) + Inches(0.3)
# Keep height 1.3, width proportional to 342/600 ratio = 0.57
photo_h = Inches(1.3)
photo_w = Inches(1.3 * 342 / 600)  # ~0.74 inches

pic = slide.shapes.add_picture(photo_path, photo_left, photo_top, photo_w, photo_h)

prs.save('AgentRep-PitchDeck.pptx')
print('Fixed olfa photo with correct aspect ratio')
