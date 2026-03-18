from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

prs = Presentation('AgentRep-PitchDeck.pptx')

# Fix ERC-8004 split text on slide 6 (index 5)
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            txt = para.text.strip()
            if txt == 'ERC-' or txt == 'ERC-\n8004' or 'ERC-\n' in repr(para.text):
                # Fix: set text to single line
                for run in para.runs:
                    if 'ERC-' in run.text:
                        run.text = 'ERC-8004'
                # Remove extra paragraphs with just "8004"
            if txt == '8004':
                para.text = ''

# Also check for runs
for shape in slide6.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        full = tf.text
        if 'ERC-' in full and '8004' in full and len(full) < 20:
            # Rebuild as single text
            paras_text = []
            for p in tf.paragraphs:
                paras_text.append(p.text.strip())
            combined = ' '.join(t for t in paras_text if t)
            if 'ERC-' in combined and '8004' in combined and 'Identity' not in combined:
                # This is the title box - set first para
                tf.paragraphs[0].text = 'ERC-8004'
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(0x00, 0xD4, 0x7E)
                # Clear remaining
                for i in range(1, len(tf.paragraphs)):
                    if tf.paragraphs[i].text.strip() == '8004':
                        tf.paragraphs[i].text = ''

prs.save('AgentRep-PitchDeck.pptx')
print('Fixed ERC-8004 split text and saved')
