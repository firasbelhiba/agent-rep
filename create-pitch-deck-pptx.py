"""
AgentRep Pitch Deck Generator using python-pptx
Creates a professional dark-themed pitch deck with accurate whitepaper content.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import subprocess

# Theme constants
BG = RGBColor(0x0d, 0x0d, 0x1a)
PRIMARY = RGBColor(0x82, 0x59, 0xef)
LIGHT_PURPLE = RGBColor(0xb4, 0x7a, 0xff)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x99, 0x99, 0xaa)
DARK_CARD = RGBColor(0x1a, 0x1a, 0x2e)
CARD_BORDER = RGBColor(0x2a, 0x2a, 0x4a)
RED_ACCENT = RGBColor(0xff, 0x44, 0x44)
GREEN = RGBColor(0x10, 0xb9, 0x81)
BLUE = RGBColor(0x3b, 0x82, 0xf6)
YELLOW = RGBColor(0xff, 0xd4, 0x3b)
ORANGE = RGBColor(0xff, 0x8c, 0x00)

TITLE_FONT = "Arial Black"
BODY_FONT = "Calibri"

LOGO_PATH = os.path.join(os.path.dirname(__file__), "public", "logo-trimmed.png")
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), "AgentRep-PitchDeck.pptx")
PDF_PATH = os.path.join(os.path.dirname(__file__), "AgentRep-PitchDeck.pdf")

TOTAL_SLIDES = 12

def inch(val):
    return Inches(val)

def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_oval(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(x), inch(y), inch(w), inch(h))
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, font_size=12, font_name=BODY_FONT,
             color=WHITE, bold=False, align="left", valign="top", margin=None):
    txBox = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    if margin is not None:
        tf.margin_left = Emu(margin)
        tf.margin_right = Emu(margin)
        tf.margin_top = Emu(margin)
        tf.margin_bottom = Emu(margin)

    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold

    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    return txBox

def add_bullet_text(slide, items, x, y, w, h, font_size=11, font_name=BODY_FONT, color=GRAY, bullet_color='b47aff'):
    txBox = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(4)
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        from lxml import etree
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '\u2022')
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
        srgbClr.set('val', bullet_color)

    return txBox

def add_footer_bar(slide):
    add_rect(slide, 0, 5.35, 10, 0.275, fill_color=PRIMARY)

def add_slide_number(slide, num):
    add_text(slide, f"{num} / {TOTAL_SLIDES}", 8.5, 5.2, 1.2, 0.3,
             font_size=9, color=GRAY, align="right")

def add_card(slide, x, y, w, h, accent_color=None, accent_side="top"):
    card = add_rect(slide, x, y, w, h, fill_color=DARK_CARD, line_color=CARD_BORDER, line_width=1)
    if accent_color:
        if accent_side == "top":
            add_rect(slide, x, y, w, 0.05, fill_color=accent_color)
        elif accent_side == "left":
            add_rect(slide, x, y, 0.05, h, fill_color=accent_color)
    return card

def add_title_section(slide, title):
    add_text(slide, title, 0.5, 0.3, 9, 0.7,
             font_size=28, font_name=TITLE_FONT, color=WHITE, bold=True)
    add_rect(slide, 0.5, 0.95, 1.5, 0.04, fill_color=PRIMARY)


def main():
    prs = Presentation()
    prs.slide_width = inch(10)
    prs.slide_height = inch(5.625)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ════════════════════════════════════════════
    # SLIDE 1: Title
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_rect(slide, 0, 0, 10, 0.04, fill_color=PRIMARY)

    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, inch(3.5), inch(0.6), inch(3), inch(0.75))

    add_text(slide, "On-Chain Reputation for AI Agents", 1, 2.0, 8, 0.6,
             font_size=22, font_name=BODY_FONT, color=LIGHT_PURPLE, bold=True, align="center")

    add_text(slide, "Built on Hedera", 1, 2.7, 8, 0.4,
             font_size=14, color=GRAY, align="center")

    add_rect(slide, 4, 3.2, 2, 0.03, fill_color=PRIMARY)

    add_text(slide, "agentrep.xyz", 1, 3.5, 8, 0.35,
             font_size=13, color=PRIMARY, align="center")

    add_footer_bar(slide)

    # ════════════════════════════════════════════
    # SLIDE 2: Team
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_title_section(slide, "The Team")

    add_card(slide, 1, 1.5, 8, 3.2, accent_color=PRIMARY, accent_side="left")

    add_text(slide, "Firas Belhiba", 1.5, 1.7, 7, 0.5,
             font_size=24, font_name=TITLE_FONT, color=WHITE, bold=True)

    add_text(slide, "Full-Stack Blockchain Developer", 1.5, 2.2, 7, 0.4,
             font_size=16, color=PRIMARY)

    add_text(slide, "Dar Blockchain", 1.5, 2.6, 7, 0.35,
             font_size=14, color=GRAY)

    add_text(slide, "Solo developer with deep experience in Hedera, Solidity, Next.js, and NestJS. Built AgentRep end-to-end as a complete decentralized reputation protocol for AI agents.",
             1.5, 3.1, 7, 1.2, font_size=13, color=GRAY)

    badges = ["Hedera", "Solidity", "Next.js", "NestJS", "TypeScript"]
    bx = 1.5
    for badge in badges:
        bw = len(badge) * 0.1 + 0.4
        add_rect(slide, bx, 4.15, bw, 0.32, fill_color=DARK_CARD, line_color=PRIMARY, line_width=0.5)
        add_text(slide, badge, bx, 4.15, bw, 0.32,
                 font_size=10, color=LIGHT_PURPLE, align="center")
        bx += bw + 0.15

    add_slide_number(slide, 2)
    add_footer_bar(slide)

    # ════════════════════════════════════════════
    # SLIDE 3: The Problem (4 problems from whitepaper)
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_title_section(slide, "The Problem")

    add_text(slide, "AI agents are proliferating, but there is no way to know which ones to trust.",
             0.5, 1.15, 9, 0.4, font_size=13, color=GRAY)

    problems = [
        {"title": "No Verifiable\nHistory", "icon": "\u2753",
         "desc": "Can't verify an agent's track record. Past performance is invisible and unverifiable."},
        {"title": "Sybil\nAttacks", "icon": "\u26d4",
         "desc": "Fake agents inflating reputation through fake feedback. No identity verification."},
        {"title": "No\nAccountability", "icon": "\u26a0",
         "desc": "Dishonest feedback with zero consequences. No staking, no slashing, no recourse."},
        {"title": "Centralized Trust\nBottlenecks", "icon": "\u26e8",
         "desc": "Centralized databases can be tampered with. Single points of failure for trust data."}
    ]

    for i, prob in enumerate(problems):
        cx = 0.35 + i * 2.4
        add_card(slide, cx, 1.7, 2.15, 3.3, accent_color=RED_ACCENT)
        add_text(slide, prob["icon"], cx, 1.95, 2.15, 0.55,
                 font_size=28, color=RGBColor(0xff, 0x6b, 0x6b), align="center", bold=True)
        add_text(slide, prob["title"], cx + 0.1, 2.6, 1.95, 0.6,
                 font_size=12, font_name=TITLE_FONT, color=WHITE, bold=True, align="center")
        add_text(slide, prob["desc"], cx + 0.1, 3.25, 1.95, 1.5,
                 font_size=10, color=GRAY, align="center")

    add_slide_number(slide, 3)
    add_footer_bar(slide)

    # ════════════════════════════════════════════
    # SLIDE 4: The Solution — The AgentRep Protocol (4 mechanisms)
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_title_section(slide, "The AgentRep Protocol")

    add_text(slide, "Four mechanisms to establish decentralized trust for AI agents",
             0.5, 1.1, 9, 0.35, font_size=13, color=GRAY)

    mechanisms = [
        {"title": "Reputation-Weighted\nFeedback",
         "desc": "weight = 0.2 + 0.8 \u00d7 (giverScore/1000)\nNew agents have low influence,\nestablished agents have high influence.",
         "color": PRIMARY},
        {"title": "Stake-Based\nAccountability",
         "desc": "5 HBAR minimum stake.\n10% slashed per upheld dispute\nvia AgentRepStaking.sol smart contract.",
         "color": GREEN},
        {"title": "Cross-Validation &\nOutlier Detection",
         "desc": "Z-score analysis. Feedback >1.5\nstd dev auto-discounted\n(down to 0.1x weight).",
         "color": BLUE},
        {"title": "Validation of\nValidators",
         "desc": "validationWeight = 0.3 + 0.7\n\u00d7 (validatorScore/1000)\nRecursive web of trust.",
         "color": ORANGE}
    ]

    for i, mech in enumerate(mechanisms):
        cx = 0.35 + i * 2.4
        add_card(slide, cx, 1.6, 2.15, 3.5, accent_color=mech["color"])
        # Mechanism number
        circle = add_oval(slide, cx + 0.75, 1.8, 0.65, 0.65, fill_color=mech["color"])
        add_text(slide, str(i + 1), cx + 0.75, 1.8, 0.65, 0.65,
                 font_size=18, font_name=TITLE_FONT, color=WHITE, bold=True, align="center")
        add_text(slide, mech["title"], cx + 0.1, 2.6, 1.95, 0.6,
                 font_size=11, font_name=TITLE_FONT, color=WHITE, bold=True, align="center")
        add_text(slide, mech["desc"], cx + 0.1, 3.25, 1.95, 1.6,
                 font_size=9, color=GRAY, align="center")

    add_slide_number(slide, 4)
    add_footer_bar(slide)

    # ════════════════════════════════════════════
    # SLIDE 5: Reputation Score Algorithm
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_title_section(slide, "Reputation Score Algorithm")

    add_text(slide, "Composite score 0\u20131000 from 4 components", 0.5, 1.1, 9, 0.35,
             font_size=13, color=GRAY)

    # Score components (left side)
    components = [
        {"name": "Quality", "weight": "30%", "max": "300",
         "desc": "Reputation-weighted + outlier-discounted feedback average",
         "color": RGBColor(0x82, 0x59, 0xef)},
        {"name": "Reliability", "weight": "30%", "max": "300",
         "desc": "Validator-weighted validation score average",
         "color": RGBColor(0x6c, 0x45, 0xd9)},
        {"name": "Activity", "weight": "20%", "max": "200",
         "desc": "Logarithmic scale of total interactions",
         "color": RGBColor(0xb4, 0x7a, 0xff)},
        {"name": "Consistency", "weight": "20%", "max": "200",
         "desc": "Low standard deviation = higher score",
         "color": RGBColor(0x99, 0x66, 0xee)}
    ]

    for i, comp in enumerate(components):
        cy = 1.6 + i * 0.82
        add_card(slide, 0.5, cy, 5.2, 0.72, accent_color=comp["color"], accent_side="left")
        add_text(slide, comp["name"], 0.7, cy + 0.05, 1.3, 0.35,
                 font_size=13, font_name=TITLE_FONT, color=WHITE, bold=True)
        add_text(slide, f"{comp['weight']} (max {comp['max']})", 2.0, cy + 0.05, 1.5, 0.35,
                 font_size=11, color=comp["color"], bold=True)
        add_text(slide, comp["desc"], 0.7, cy + 0.38, 4.8, 0.3,
                 font_size=10, color=GRAY)

    # Trust Tiers (right side)
    add_card(slide, 6.0, 1.6, 3.5, 3.5, accent_color=PRIMARY)
    add_text(slide, "Trust Tiers", 6.2, 1.75, 3.1, 0.35,
             font_size=14, font_name=TITLE_FONT, color=WHITE, bold=True, align="center")

    tiers = [
        {"name": "UNVERIFIED", "range": "< 200", "color": GRAY},
        {"name": "VERIFIED", "range": "\u2265 200, \u2265 3 activities", "color": BLUE},
        {"name": "TRUSTED", "range": "\u2265 500, \u2265 10 activities", "color": GREEN},
        {"name": "ELITE", "range": "\u2265 800, \u2265 20 activities", "color": RGBColor(0xff, 0xd4, 0x3b)}
    ]

    for i, tier in enumerate(tiers):
        ty = 2.25 + i * 0.65
        add_rect(slide, 6.3, ty, 3.0, 0.5, fill_color=BG, line_color=tier["color"], line_width=1)
        add_text(slide, tier["name"], 6.4, ty + 0.03, 1.5, 0.25,
                 font_size=11, font_name=TITLE_FONT, color=tier["color"], bold=True)
        add_text(slide, tier["range"], 6.4, ty + 0.26, 2.8, 0.2,
                 font_size=9, color=GRAY)

    add_slide_number(slide, 5)
    add_footer_bar(slide)

    # ════════════════════════════════════════════
    # SLIDE 6: On-Chain Architecture
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_title_section(slide, "On-Chain Architecture")

    # 3 HCS Topics (left)
    add_text(slide, "HCS Topics", 0.5, 1.15, 3, 0.35,
             font_size=14, font_name=TITLE_FONT, color=LIGHT_PURPLE)

    topics = [
        {"name": "Identity Topic", "events": "AGENT_REGISTERED, URI_UPDATED"},
        {"name": "Feedback Topic", "events": "FEEDBACK_SUBMITTED, STAKE_DEPOSITED, STAKE_SLASHED"},
        {"name": "Validation Topic", "events": "VALIDATION_REQUESTED, VALIDATION_RESPONDED"}
    ]

    for i, topic in enumerate(topics):
        ty = 1.6 + i * 0.85
        add_card(slide, 0.5, ty, 4.5, 0.75, accent_color=PRIMARY, accent_side="left")
        add_text(slide, topic["name"], 0.7, ty + 0.05, 2.0, 0.3,
                 font_size=12, font_name=TITLE_FONT, color=WHITE, bold=True)
        add_text(slide, topic["events"], 0.7, ty + 0.38, 4.1, 0.3,
                 font_size=9, color=GRAY)

    # Smart Contract (right)
    add_text(slide, "Smart Contract", 5.5, 1.15, 4, 0.35,
             font_size=14, font_name=TITLE_FONT, color=LIGHT_PURPLE)

    add_card(slide, 5.5, 1.6, 4.0, 1.2, accent_color=GREEN)
    add_text(slide, "AgentRepStaking.sol", 5.7, 1.72, 3.6, 0.3,
             font_size=12, font_name=TITLE_FONT, color=WHITE, bold=True)
    add_text(slide, "Contract: 0.0.8264743", 5.7, 2.02, 3.6, 0.25,
             font_size=10, color=GREEN)
    add_text(slide, "Functions: stake, slash, unstake, getStake", 5.7, 2.3, 3.6, 0.25,
             font_size=9, color=GRAY)

    # Standards (bottom)
    add_text(slide, "Standards", 0.5, 4.1, 9, 0.35,
             font_size=14, font_name=TITLE_FONT, color=LIGHT_PURPLE)

    standards = [
        {"name": "ERC-8004", "desc": "Identity + Reputation + Validation"},
        {"name": "HCS-10", "desc": "Agent Communication"},
        {"name": "HCS-11", "desc": "Agent Profiles"}
    ]

    for i, std in enumerate(standards):
        sx = 0.5 + i * 3.15
        add_rect(slide, sx, 4.5, 2.85, 0.6, fill_color=DARK_CARD, line_color=CARD_BORDER, line_width=1)
        add_text(slide, std["name"], sx + 0.1, 4.52, 1.0, 0.25,
                 font_size=11, font_name=TITLE_FONT, color=PRIMARY, bold=True)
        add_text(slide, std["desc"], sx + 0.1, 4.78, 2.6, 0.25,
                 font_size=9, color=GRAY)

    add_slide_number(slide, 6)
    add_footer_bar(slide)

    # ════════════════════════════════════════════
    # SLIDE 7: Key Features (6 features)
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_title_section(slide, "Key Features")

    features = [
        {"title": "Reputation Scoring", "desc": "Weighted 0\u20131000 score with reputation-weighted feedback and outlier detection"},
        {"title": "Staking & Slashing", "desc": "HBAR staking via smart contract. 5 HBAR min, 10% slash per upheld dispute"},
        {"title": "HCS-10 Messaging", "desc": "Agent-to-agent communication with full connection lifecycle management"},
        {"title": "Community Reviews", "desc": "Wallet-verified feedback via HashPack. Real users, not anonymous ratings"},
        {"title": "Dispute Resolution", "desc": "On-chain arbitration with designated arbiter. Stake slashing for bad actors"},
        {"title": "Developer SDK", "desc": "npm: agent-rep-sdk. TypeScript SDK with AgentRunner for autonomous agents"}
    ]

    for row in range(2):
        for col in range(3):
            idx = row * 3 + col
            cx = 0.5 + col * 3.1
            cy = 1.15 + row * 2.15

            add_card(slide, cx, cy, 2.8, 1.9, accent_color=PRIMARY, accent_side="left")

            add_text(slide, features[idx]["title"], cx + 0.2, cy + 0.2, 2.4, 0.45,
                     font_size=13, font_name=TITLE_FONT, color=WHITE, bold=True)
            add_text(slide, features[idx]["desc"], cx + 0.2, cy + 0.75, 2.4, 0.95,
                     font_size=10, color=GRAY)

    add_slide_number(slide, 7)
    add_footer_bar(slide)

    # ════════════════════════════════════════════
    # SLIDE 8: Security Model
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_title_section(slide, "Security Model")

    add_text(slide, "Multi-layered defense against reputation gaming",
             0.5, 1.1, 9, 0.35, font_size=13, color=GRAY)

    security = [
        {"attack": "Sybil Attacks",
         "defense": "HCS-10 connection required + stake cost",
         "icon": "\u26d4", "color": RED_ACCENT},
        {"attack": "Reputation Manipulation",
         "defense": "Reputation-weighted feedback + outlier detection (Z-score, 1.5 std dev threshold)",
         "icon": "\u26e8", "color": PRIMARY},
        {"attack": "Feedback Spam",
         "defense": "1 HBAR stake required + rate limiting (20 requests/hour)",
         "icon": "\u26a0", "color": ORANGE},
        {"attack": "Data Tampering",
         "defense": "All events logged on HCS \u2014 immutable, verifiable on HashScan",
         "icon": "\u2713", "color": GREEN}
    ]

    for i, sec in enumerate(security):
        cy = 1.6 + i * 0.9
        add_card(slide, 0.5, cy, 9, 0.8, accent_color=sec["color"], accent_side="left")

        add_text(slide, sec["icon"], 0.7, cy + 0.1, 0.5, 0.5,
                 font_size=22, color=sec["color"], align="center", bold=True)

        add_text(slide, sec["attack"], 1.3, cy + 0.08, 2.5, 0.35,
                 font_size=13, font_name=TITLE_FONT, color=WHITE, bold=True)

        add_text(slide, "\u2192", 3.8, cy + 0.08, 0.4, 0.35,
                 font_size=16, color=sec["color"], align="center", bold=True)

        add_text(slide, sec["defense"], 4.3, cy + 0.08, 5.0, 0.6,
                 font_size=11, color=GRAY)

    add_slide_number(slide, 8)
    add_footer_bar(slide)

    # ════════════════════════════════════════════
    # SLIDE 9: Deployed Resources
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_title_section(slide, "Deployed Resources")

    add_text(slide, "All deployed and verifiable on Hedera Testnet",
             0.5, 1.15, 9, 0.4, font_size=13, color=GRAY)

    on_chain = [
        {"label": "Smart Contract", "value": "0.0.8264743", "note": "Source-verified on Sourcify"},
        {"label": "Identity Topic", "value": "0.0.8264956", "note": "Agent registration records"},
        {"label": "Feedback Topic", "value": "0.0.8264959", "note": "Feedback & staking events"},
        {"label": "Validation Topic", "value": "0.0.8264962", "note": "Validation requests & responses"}
    ]

    for i, item in enumerate(on_chain):
        cy = 1.7 + i * 0.85
        add_rect(slide, 0.5, cy, 5.5, 0.7, fill_color=DARK_CARD, line_color=CARD_BORDER, line_width=1)

        add_text(slide, "\u2713", 0.65, cy + 0.1, 0.4, 0.5,
                 font_size=18, color=RGBColor(0x51, 0xcf, 0x66), bold=True)
        add_text(slide, item["label"], 1.1, cy + 0.05, 1.8, 0.35,
                 font_size=12, font_name=TITLE_FONT, color=WHITE, bold=True)
        add_text(slide, item["value"], 3.0, cy + 0.05, 1.6, 0.35,
                 font_size=12, color=PRIMARY)
        add_text(slide, item["note"], 1.1, cy + 0.38, 4.5, 0.3,
                 font_size=9, color=GRAY)

    # Tech stack
    add_text(slide, "Tech Stack", 6.5, 1.7, 3, 0.4,
             font_size=14, font_name=TITLE_FONT, color=LIGHT_PURPLE)

    tech = ["Next.js", "NestJS", "TypeORM", "PostgreSQL", "Solidity", "Hedera SDK"]
    for i, t in enumerate(tech):
        cy = 2.2 + i * 0.48
        add_rect(slide, 6.5, cy, 2.8, 0.38, fill_color=DARK_CARD, line_color=CARD_BORDER, line_width=0.5)
        add_text(slide, "\u2713", 6.6, cy, 0.3, 0.38,
                 font_size=14, color=RGBColor(0x51, 0xcf, 0x66), bold=True)
        add_text(slide, t, 7.0, cy, 2.2, 0.38,
                 font_size=11, color=WHITE)

    add_slide_number(slide, 9)
    add_footer_bar(slide)

    # ════════════════════════════════════════════
    # SLIDE 10: Roadmap & Key Learnings
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_title_section(slide, "Roadmap & Key Learnings")

    # Roadmap (left)
    add_card(slide, 0.5, 1.2, 5.3, 4.0, accent_color=PRIMARY)
    add_text(slide, "Roadmap", 0.7, 1.35, 2.5, 0.35,
             font_size=14, font_name=TITLE_FONT, color=WHITE, bold=True)

    # Phase 1
    add_text(slide, "\u2713 Phase 1 (Complete)", 0.8, 1.75, 4.8, 0.3,
             font_size=11, color=GREEN, bold=True)
    add_text(slide, "ERC-8004 registries, HCS-10/11 integration, on-chain logging",
             0.8, 2.0, 4.8, 0.3, font_size=9, color=GRAY)

    # Phase 2
    add_text(slide, "\u2713 Phase 2 (Complete)", 0.8, 2.35, 4.8, 0.3,
             font_size=11, color=GREEN, bold=True)
    add_text(slide, "Reputation-weighted feedback, staking/slashing, outlier detection, SDK, smart contract",
             0.8, 2.6, 4.8, 0.35, font_size=9, color=GRAY)

    # Phase 3
    add_text(slide, "\u2192 Phase 3 (Planned)", 0.8, 3.0, 4.8, 0.3,
             font_size=11, color=LIGHT_PURPLE, bold=True)
    phase3_items = [
        "DAO governance for dispute resolution",
        "Reputation decay over time",
        "Cross-chain bridging",
        "Mainnet deployment",
        "AI arbiter agents"
    ]
    add_bullet_text(slide, phase3_items, 0.8, 3.3, 4.8, 1.8, font_size=9, color=GRAY)

    # Key Learnings (right)
    add_card(slide, 6.1, 1.2, 3.4, 4.0, accent_color=YELLOW)
    add_text(slide, "Key Learnings", 6.3, 1.35, 3.0, 0.35,
             font_size=14, font_name=TITLE_FONT, color=WHITE, bold=True)

    learnings = [
        "HCS-10 is powerful for structured agent-to-agent communication",
        "Reputation algorithms need careful weighting to prevent gaming",
        "Economic incentives (staking) are critical for establishing trust"
    ]
    add_bullet_text(slide, learnings, 6.3, 1.85, 3.0, 3.0, font_size=10, color=GRAY, bullet_color='ffd43b')

    add_slide_number(slide, 10)
    add_footer_bar(slide)

    # ════════════════════════════════════════════
    # SLIDE 11: Live Demo
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_title_section(slide, "Live Demo")

    # Demo highlights card (left)
    add_card(slide, 0.5, 1.3, 3.5, 3.8, accent_color=PRIMARY, accent_side="left")
    add_text(slide, "Demo Highlights", 0.7, 1.45, 3.1, 0.4,
             font_size=12, font_name=TITLE_FONT, color=LIGHT_PURPLE, bold=True)

    highlights = [
        "Agent registration with HBAR",
        "Reputation scoring in action",
        "HCS-10 agent messaging",
        "Feedback / validation / dispute flow",
        "SDK-powered agent listener"
    ]
    add_bullet_text(slide, highlights, 0.7, 1.9, 3.1, 3.0, font_size=10, color=GRAY)

    # Play button area (center-right)
    outer_circle = add_oval(slide, 5.2, 1.4, 2.4, 2.4, line_color=PRIMARY, line_width=2)
    inner_circle = add_oval(slide, 5.6, 1.8, 1.6, 1.6, fill_color=PRIMARY)
    add_text(slide, "\u25b6", 5.8, 2.0, 1.2, 1.2,
             font_size=36, color=WHITE, align="center", bold=True)

    add_text(slide, "Watch the Demo", 4.5, 3.9, 4.5, 0.5,
             font_size=18, font_name=TITLE_FONT, color=WHITE, bold=True, align="center")

    add_text(slide, "[INSERT YOUTUBE LINK]", 4.5, 4.3, 4.5, 0.3,
             font_size=12, color=PRIMARY, align="center")

    add_text(slide, "Live on Hedera Testnet  \u2014  agentrep.xyz", 4.5, 4.6, 4.5, 0.3,
             font_size=12, color=GRAY, align="center")

    add_slide_number(slide, 11)
    add_footer_bar(slide)

    # ════════════════════════════════════════════
    # SLIDE 12: Thank You
    # ════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_rect(slide, 0, 0, 10, 0.04, fill_color=PRIMARY)

    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, inch(3.75), inch(0.5), inch(2.5), inch(0.625))

    add_text(slide, "AgentRep", 1, 1.6, 8, 0.6,
             font_size=24, font_name=TITLE_FONT, color=WHITE, bold=True, align="center")

    add_text(slide, "agentrep.xyz", 1, 2.2, 8, 0.4,
             font_size=16, color=PRIMARY, align="center")

    add_text(slide, "Built on Hedera", 1, 2.6, 8, 0.35,
             font_size=13, color=GRAY, align="center")

    add_rect(slide, 3.5, 3.1, 3, 0.03, fill_color=PRIMARY)

    # Links card
    add_rect(slide, 2.5, 3.3, 5, 1.2, fill_color=DARK_CARD, line_color=CARD_BORDER, line_width=1)

    add_text(slide, "GitHub:  github.com/firasbelhiba/agent-rep", 2.8, 3.45, 4.5, 0.3,
             font_size=11, color=GRAY)

    add_text(slide, "npm:     agent-rep-sdk", 2.8, 3.8, 4.5, 0.3,
             font_size=11, color=GRAY)

    add_text(slide, "Web:     agentrep.xyz", 2.8, 4.1, 4.5, 0.3,
             font_size=11, color=GRAY)

    add_footer_bar(slide)

    # ── Save ──
    prs.save(OUTPUT_PATH)
    print(f"PPTX created: {OUTPUT_PATH}")

    # ── Convert to PDF ──
    vbs_script = os.path.join(os.path.dirname(__file__), "convert_to_pdf.vbs")
    with open(vbs_script, 'w') as f:
        f.write(f'''Dim pptApp, pptPres
Set pptApp = CreateObject("PowerPoint.Application")
pptApp.Visible = True
On Error Resume Next
Set pptPres = pptApp.Presentations.Open("{OUTPUT_PATH}", -1, 0, -1)
If Err.Number <> 0 Then
    WScript.Echo "Error opening: " & Err.Description
    pptApp.Quit
    WScript.Quit 1
End If
On Error GoTo 0
pptPres.SaveAs "{PDF_PATH}", 32
pptPres.Close
pptApp.Quit
WScript.Echo "PDF created successfully"
''')

    result = subprocess.run(["cscript", "//nologo", vbs_script],
                          capture_output=True, text=True, timeout=60)
    print("STDOUT:", result.stdout.strip())
    if result.stderr:
        print("STDERR:", result.stderr.strip())
    if result.returncode != 0:
        print(f"PDF conversion failed with code {result.returncode}")
    else:
        print(f"PDF created: {PDF_PATH}")

    # Cleanup temp files
    try:
        os.remove(vbs_script)
    except:
        pass


if __name__ == "__main__":
    main()
